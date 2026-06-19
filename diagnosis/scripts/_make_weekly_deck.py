"""Build the plain weekly-report deck (white slides, bullets + figures).
Revised 2026-06-19: plainer language + per-slide 'In plain terms' takeaway."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from pathlib import Path

FIG = Path("results/figures/week_report")
import os
_pref = Path("../Weekly_Report_2026-06-19.pptx")
# If the preferred file is locked (open in PowerPoint), fall back to a _v2 name.
OUT = _pref
try:
    if _pref.exists():
        with open(_pref, "ab"):
            pass
except PermissionError:
    OUT = Path("../Weekly_Report_2026-06-19_v2.pptx")
DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x55, 0x55, 0x55)
ACC = RGBColor(0x1B, 0x5E, 0x20)   # deep green for the plain-terms line

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]


def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = DARK
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(14); r2.font.italic = True; r2.font.color.rgb = GREY


def add_takeaway(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.18), Inches(12.1), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "In plain terms:  "; r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = ACC
    r2 = p.add_run(); r2.text = text; r2.font.size = Pt(15); r2.font.italic = True; r2.font.color.rgb = ACC


def add_bullets(slide, bullets, x=0.6, y=1.95, w=6.5, h=5.0, size=15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lvl, txt, *bold) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(7)
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl); r.font.color.rgb = DARK
        if bold and bold[0]:
            r.font.bold = True


def add_image_fit(slide, path, x, y, maxw, maxh):
    w, h = Image.open(path).size; ar = w / h
    bw, bh = maxw, maxw / ar
    if bh > maxh:
        bh, bw = maxh, maxh * ar
    slide.shapes.add_picture(str(path), Inches(x + (maxw - bw) / 2),
                             Inches(y + (maxh - bh) / 2), Inches(bw), Inches(bh))


# 1 — title
s = prs.slides.add_slide(blank)
add_title(s, "Boundary-Blind World Models — Weekly Progress", "Week of 2026-06-12 → 06-19")
add_bullets(s, [
    (0, "Goal: make these AI 'world models' plan contact tasks (push, pick-and-place) — today they fail.", True),
    (0, "This week we found exactly why: the model predicts almost the same result no matter which action you try."),
    (0, "Its 'eyes' and 'memory' are fine — the earlier 'fixes' only looked like fixes."),
    (0, "On the one task it does handle (reaching), we already beat the published result: 94% vs 44.8%."),
], y=2.0, w=12.2, h=4.8, size=18)

# 2 — recap + question
s = prs.slides.add_slide(blank)
add_title(s, "Recap & the question")
add_takeaway(s, "the arm reaches the target, but the robot never actually moves the object.")
add_bullets(s, [
    (0, "A 'world model' imagines what the robot will see after an action; the planner uses it to pick actions."),
    (0, "Reaching an empty-space target works (94%). The moment the arm must move an OBJECT, it fails (0%)."),
    (0, "Question: is the problem the planner, the model's vision, or the model's predictions?"),
])
add_image_fit(s, FIG / "closed_loop.png", 7.3, 2.0, 5.7, 4.8)

# 3 — three fixes + the catch
s = prs.slides.add_slide(blank)
add_title(s, "Three fixes — and why we dug deeper")
add_takeaway(s, "two fixes looked promising, so we built a sharper ruler to check if they were real.")
add_bullets(s, [
    (0, "Fix 1 — score actions by how much they move the object: tiny gain."),
    (0, "Fix 2 — push the planner to seek contact: no help (slightly worse)."),
    (0, "Fix 3 — correct the model's predictions: looked like the best..."),
    (0, "...but all three left success at 0/16. So we built a more precise way to read the object.", True),
])
add_image_fit(s, FIG / "fix_ladder.png", 6.9, 2.2, 6.1, 4.4)

# 4 — not vision, not memory
s = prs.slides.add_slide(blank)
add_title(s, "It's not the model's vision or its memory")
add_takeaway(s, "the model's 'eyes' (where the object is) and 'memory' (what just happened) are both fine.")
add_bullets(s, [
    (0, "We read the object's position out of the model two ways; the old readout was blurry (6 cm)."),
    (0, "A better readout pins it to 2 cm — inside the 5 cm success zone 92% of the time. The info was always there.", True),
    (0, "The model also correctly replays what really happened (~3 cm over 6 steps)."),
    (0, "And Fix 3 is no better than no-fix under the precise readout — it had just gamed the blurry one.", True),
], y=1.9, w=12.3, h=2.0)
add_image_fit(s, FIG / "precision_pooled_vs_spatial.png", 2.5, 4.0, 8.3, 3.3)

# 5 — the real problem
s = prs.slides.add_slide(blank)
add_title(s, "The real problem: it can't tell actions apart")
add_takeaway(s, "the model imagines almost the same object outcome for every action — so the planner can't choose.")
add_bullets(s, [
    (0, "Planning = imagine many actions, keep the one that best moves the object to the goal."),
    (0, "In reality, different actions send the object ~2.4 cm apart."),
    (0, "But the model imagines them all landing in nearly the same spot (0.74 cm) — unrelated to the truth.", True),
    (0, "So every action looks equally good to the planner. That is the core failure.", True),
])
add_image_fit(s, FIG / "counterfactual_smear.png", 7.5, 2.0, 5.5, 4.7)

# 6 — resolved + next
s = prs.slides.add_slide(blank)
add_title(s, "Where it stands & what's next")
add_takeaway(s, "vision ✓, memory ✓, imagining different actions ✗ — that last part is what we now fix.")
add_bullets(s, [
    (0, "The failure, located:", True),
    (1, "Vision (where the object is): FINE — 2 cm."),
    (1, "Memory (replaying what happened): FINE — ~3 cm over 6 steps."),
    (1, "Imagining different actions differently: BROKEN — the culprit."),
    (0, "Fix 3 was a mirage (it gamed the blurry readout) — an honest negative."),
    (0, "Running now: give the planner an accurate object target + the one part that DOES tell actions"),
    (1, "apart, weighted to drive the choice. Works → first contact successes; if not → strengthen that part."),
    (0, "Reaching already beats the paper (94% vs 44.8%).", True),
], y=1.95, w=12.3)

prs.save(str(OUT))
print("wrote", OUT.resolve())
